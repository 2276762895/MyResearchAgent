import streamlit as st
import os
# --- 【新增】设置国内镜像 ---
os.environ['HF_ENDPOINT'] = 'https://hf-mirror.com'
import tempfile
from langchain_community.document_loaders import PyPDFLoader
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_community.vectorstores import FAISS
from langchain_openai import ChatOpenAI
from langchain.chains import RetrievalQA

import json
from pptx import Presentation


# --- ✨ 新增功能：PPT 生成函数 ---
def generate_ppt_file(topic, content_json):
    """
    输入：PPT主题，和 DeepSeek 生成的 JSON 内容
    输出：生成的 PPT 文件路径
    """
    try:
        # 1. 创建 PPT 对象
        prs = Presentation()

        # 2. 解析 JSON 数据
        # 有时候大模型会包裹 markdown 代码块，需要清洗
        clean_json = content_json.replace("```json", "").replace("```", "").strip()
        data = json.loads(clean_json)

        # 3. 生成 封面页
        slide_layout = prs.slide_layouts[0]  # 0 是标题页
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = topic
        slide.placeholders[1].text = "Powered by DeepSeek & Python"

        # 4. 循环生成 正文页
        for page in data['pages']:
            slide_layout = prs.slide_layouts[1]  # 1 是标题+内容页
            slide = prs.slides.add_slide(slide_layout)

            # 填标题
            slide.shapes.title.text = page['title']

            # 填内容（把列表变成带点的文本）
            tf = slide.placeholders[1].text_frame
            for point in page['content']:
                p = tf.add_paragraph()
                p.text = point
                p.level = 0

        # 5. 保存文件
        output_path = "generated_ppt.pptx"
        prs.save(output_path)
        return output_path

    except Exception as e:
        st.error(f"PPT生成失败，请重试。错误信息: {e}")
        return None

# --- 页面设置 ---
st.set_page_config(page_title="科研论文智能助手", layout="wide")
st.title("🎓 课题组科研助手 - 文献阅读版")

# --- 侧边栏：设置与文件上传 ---
with st.sidebar:
    st.header("⚙️ 设置")
    # 输入 DeepSeek API Key
    api_key = st.text_input("请输入 DeepSeek API Key", type="password")
    st.markdown("[点击申请 DeepSeek API](https://platform.deepseek.com/)")

    st.divider()

    st.header("📂 上传文献")
    uploaded_file = st.file_uploader("上传PDF文件", type=["pdf"])
    st.divider()
    st.header("📊 生成汇报PPT")
    ppt_topic = st.text_input("请输入PPT主题", value="论文汇报")

    if st.button("开始生成 PPT"):
        if not uploaded_file:
            st.warning("请先上传论文PDF！")
        elif not api_key:
            st.warning("请填写 API Key！")
        else:
            with st.spinner("DeepSeek 正在构思 PPT 大纲..."):
                # 1. 让 DeepSeek 生成 JSON 数据
                # 这里我们利用已经存在的 qa_chain (注意：需要把 qa_chain 设为全局或 session_state，或者这里临时重新定义)
                # 为了简单稳妥，我们直接复用 process_pdf 返回的 chain
                if 'qa_chain' not in st.session_state:
                    # 如果用户还没问过问题，链可能没存，这里我们得从 process_pdf 再拿一次
                    # 为了代码简洁，建议你先在主逻辑里把 qa_chain 存进 st.session_state
                    st.warning("请先在右侧主界面等待论文读取完成！")
                else:
                    ppt_prompt = f"""
                        请根据这篇论文的内容，为主题“{ppt_topic}”生成一个PPT大纲。
                        要求：
                        1. 返回纯 JSON 格式，不要包含任何其他废话。
                        2. JSON 格式必须如下：
                        {{
                            "pages": [
                                {{"title": "研究背景", "content": ["点1", "点2"]}},
                                {{"title": "核心方法", "content": ["点1", "点2"]}},
                                {{"title": "实验结果", "content": ["点1", "点2"]}},
                                {{"title": "结论", "content": ["点1", "点2"]}}
                            ]
                        }}
                        3. 至少生成 5 页 PPT。
                        """

                    # 调用大模型
                    response = st.session_state.qa_chain.invoke({"query": ppt_prompt})
                    result_text = response["result"]

                    # 2. 调用 Python 画图
                    ppt_path = generate_ppt_file(ppt_topic, result_text)

                    if ppt_path:
                        st.success("🎉 PPT 生成成功！")
                        with open(ppt_path, "rb") as f:
                            st.download_button(
                                label="📥 点击下载 PPT",
                                data=f,
                                file_name=f"{ppt_topic}.pptx",
                                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                            )


# --- 核心逻辑函数 ---

@st.cache_resource
def process_pdf(file, api_key):
    if not file or not api_key:
        return None

    print("1. [开始] 正在保存临时文件...")
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp_file:
        tmp_file.write(file.getvalue())
        tmp_path = tmp_file.name

    print(f"2. [加载] 正在读取PDF: {tmp_path} ...")
    loader = PyPDFLoader(tmp_path)
    docs = loader.load()
    print(f"   -> PDF读取完成，共 {len(docs)} 页")

    print("3. [切分] 正在切分文本...")
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
    splits = text_splitter.split_documents(docs)
    print(f"   -> 切分完成，共生成 {len(splits)} 个文本块")

    print("4. [模型] 正在加载 Embedding 模型 (这一步最容易卡)...")
    try:
        embeddings = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
        print("   -> 模型加载成功！")
    except Exception as e:
        print(f"   -> ❌ 模型加载失败: {e}")
        raise e

    print("5. [存储] 正在写入向量数据库 (FAISS)...")
    try:
        # 使用 FAISS 替代 Chroma，不需要 SQLite 支持，不会闪退
        vectorstore = FAISS.from_documents(documents=splits, embedding=embeddings)
        print("   -> 数据库写入成功！(已切换为 FAISS)")
    except Exception as e:
        print(f"   -> ❌ FAISS 写入失败: {e}")
        raise e

    print("6. [连接] 正在初始化 DeepSeek...")
    llm = ChatOpenAI(
        model_name="deepseek-chat",
        openai_api_key=api_key,
        openai_api_base="https://api.deepseek.com",
        temperature=0.3
    )

    print("7. [完成] 准备就绪！")

    qa_chain = RetrievalQA.from_chain_type(
        llm=llm,
        retriever=vectorstore.as_retriever(search_kwargs={"k": 3}),
        return_source_documents=True
    )

    os.remove(tmp_path)
    return qa_chain

# --- 主界面逻辑 ---

if uploaded_file and api_key:
    with st.spinner("正在阅读论文，请稍候... (第一次加载模型可能需要1分钟)"):
        # 处理PDF
        qa_chain = process_pdf(uploaded_file, api_key)
        # 【新增】把这个工具存到 session_state 里，这样侧边栏也能用
        st.session_state.qa_chain = qa_chain

    st.success("✅ 论文已读取，快来问我问题吧！")

    # 初始化聊天历史
    if "messages" not in st.session_state:
        st.session_state.messages = []

    # 显示历史消息
    for message in st.session_state.messages:
        with st.chat_message(message["role"]):
            st.markdown(message["content"])

    # 获取用户输入
    if prompt := st.chat_input("这篇论文的主要贡献是什么？"):
        # 1. 显示用户问题
        st.session_state.messages.append({"role": "user", "content": prompt})
        with st.chat_message("user"):
            st.markdown(prompt)

        # 2. 调用模型回答
        with st.chat_message("assistant"):
            with st.spinner("思考中..."):
                response = qa_chain.invoke({"query": prompt})
                answer = response["result"]
                source_docs = response["source_documents"]

                # 拼接引用来源（可选）
                source_text = "\n\n> **参考片段：**\n"
                for i, doc in enumerate(source_docs):
                    source_text += f"> {i + 1}. Page {doc.metadata['page']}: {doc.page_content[:100]}...\n"

                full_response = answer + source_text
                st.markdown(full_response)

        # 3. 保存助手回答
        st.session_state.messages.append({"role": "assistant", "content": full_response})
        # 只有当有对话记录时才显示下载按钮
    if "messages" in st.session_state and len(st.session_state.messages) > 0:
        st.divider()  # 画一条分割线

        # 把对话记录转换成字符串
        chat_history_text = ""
        for msg in st.session_state.messages:
            role = "我" if msg["role"] == "user" else "AI助手"
            chat_history_text += f"[{role}]: {msg['content']}\n\n"

        # 下载按钮
        st.download_button(
            label="💾 导出对话记录 (保存为TXT)",
            data=chat_history_text,
            file_name="论文阅读记录.txt",
            mime="text/plain"
        )

else:
    st.info("👈 请在左侧输入API Key并上传PDF文件开始。")
